import subprocess
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

PPT_FILE = "자전거_안전_교육.pptx"

@pytest.fixture(scope="session", autouse=True)
def generated_ppt():
    """테스트 시작 전 PPT를 생성한다."""
    subprocess.run(["python", "make_ppt.py"], check=True)
    return Presentation(PPT_FILE)


def test_slide_count(generated_ppt):
    assert len(generated_ppt.slides) == 18


def test_slide_titles(generated_ppt):
    """각 슬라이드의 헤더 바 텍스트(첫 번째 텍스트 박스)를 확인한다."""
    expected_titles = [
        "자전거 안전 교육",           # 1 표지
        "오늘의 학습 목표",            # 2
        "자전거, 얼마나 위험할까?",     # 3
        "헬멧 착용이 왜 중요한가",      # 4
        "올바른 헬멧 착용법",           # 5
        "출발 전 점검 ABC",            # 6
        "도로 주행 규칙",              # 7
        "야간·우천 시 안전 수칙",       # 8
        "핸드폰·이어폰 사용 금지",      # 9
        "흔한 사고 유형 TOP 3",        # 10
        "사례 1: 교차로 사고",         # 11
        "사례 2: 보행자 충돌",         # 12
        "사례 3: 내리막 과속",         # 13
        "사고 발생 시 행동 요령",       # 14
        "안전 수칙 한눈에 보기",        # 15
        "OX 퀴즈",                    # 16
        "오늘의 약속",                 # 17
        "참고 자료",                   # 18
    ]
    for i, slide in enumerate(generated_ppt.slides):
        # 첫 번째 텍스트 박스 텍스트 수집
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text.strip())
        assert expected_titles[i] in texts, (
            f"슬라이드 {i+1}에서 '{expected_titles[i]}' 를 찾을 수 없음. "
            f"발견된 텍스트: {texts}"
        )
