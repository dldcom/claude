from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 색상 상수 ───────────────────────────────────────────
BLUE   = RGBColor(0x1E, 0x88, 0xE5)   # 주 색상
GREEN  = RGBColor(0x43, 0xA0, 0x47)   # 보조 색상
ORANGE = RGBColor(0xFB, 0x8C, 0x00)   # 강조 (경고)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x21, 0x21, 0x21)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)   # 연한 회색 배경

FONT_KR = "맑은 고딕"

# ─── 슬라이드 크기 (16:9 와이드) ──────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    """완전히 빈 슬라이드를 추가하고 반환한다."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_color, border_color=None):
    """채워진 직사각형 도형을 추가하고 반환한다."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font=FONT_KR):
    """텍스트 박스를 추가하고 반환한다."""
    tx_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx_box


def header_bar(slide, title_text, bg_color=BLUE):
    """슬라이드 상단에 색상 헤더 바와 제목 텍스트를 추가한다."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), bg_color)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.15),
             Inches(12.5), Inches(0.85),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def add_slide_01(prs):
    """슬라이드 1: 표지"""
    sl = blank_slide(prs)
    # 전체 배경 파란색
    add_rect(sl, Inches(0), Inches(0), W, H, BLUE)
    # 중앙 흰색 패널
    add_rect(sl, Inches(1.5), Inches(1.8), Inches(10.33), Inches(3.8), WHITE)
    # 메인 제목
    add_text(sl, "🚲 자전거 안전 교육",
             Inches(1.8), Inches(2.0), Inches(9.8), Inches(1.4),
             font_size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    # 부제목
    add_text(sl, "초등학교 5학년",
             Inches(1.8), Inches(3.5), Inches(9.8), Inches(0.6),
             font_size=22, color=DARK, align=PP_ALIGN.CENTER)
    # 날짜
    add_text(sl, "도로교통공단 교육자료 기반",
             Inches(1.8), Inches(4.1), Inches(9.8), Inches(0.5),
             font_size=16, color=RGBColor(0x75, 0x75, 0x75), align=PP_ALIGN.CENTER)
    # 하단 띠
    add_rect(sl, Inches(0), Inches(6.8), W, Inches(0.7), GREEN)
    add_text(sl, "안전한 자전거 생활을 위한 필수 교육",
             Inches(0), Inches(6.82), W, Inches(0.55),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)


def add_slide_02(prs):
    """슬라이드 2: 오늘의 학습 목표"""
    sl = blank_slide(prs)
    header_bar(sl, "오늘의 학습 목표")
    goals = [
        ("①", "자전거 기본 안전 수칙을 이해한다."),
        ("②", "주요 사고 유형과 원인을 파악한다."),
        ("③", "사고를 예방하는 방법을 익힌다."),
    ]
    colors = [BLUE, GREEN, ORANGE]
    for i, (num, text) in enumerate(goals):
        y = Inches(1.6 + i * 1.7)
        add_rect(sl, Inches(0.5), y, Inches(0.8), Inches(1.2), colors[i])
        add_text(sl, num,
                 Inches(0.5), y + Inches(0.25),
                 Inches(0.8), Inches(0.7),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(1.5), y, Inches(11.0), Inches(1.2), LGRAY)
        add_text(sl, text,
                 Inches(1.7), y + Inches(0.2),
                 Inches(10.6), Inches(0.85),
                 font_size=22, color=DARK)


def add_slide_03(prs):
    """슬라이드 3: 자전거, 얼마나 위험할까? (통계)"""
    sl = blank_slide(prs)
    header_bar(sl, "자전거, 얼마나 위험할까?", bg_color=ORANGE)
    # 통계 수치 카드 3개
    stats = [
        ("연간 사고 건수", "약 15,000건", BLUE),
        ("사망·중상 비율", "전체의 30%", ORANGE),
        ("미착용 헬멧\n사고 비율", "68%", GREEN),
    ]
    for i, (label, value, color) in enumerate(stats):
        x = Inches(0.5 + i * 4.3)
        add_rect(sl, x, Inches(1.5), Inches(3.8), Inches(2.5), color)
        add_text(sl, value,
                 x, Inches(1.7), Inches(3.8), Inches(1.2),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, label,
                 x, Inches(2.9), Inches(3.8), Inches(0.8),
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "※ 출처: 도로교통공단 교통사고분석시스템(TAAS) 2023",
             Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
             font_size=12, color=RGBColor(0x75, 0x75, 0x75))
    add_text(sl, "자전거 사고는 생각보다 훨씬 많이 발생합니다.\n헬멧만 써도 부상을 크게 줄일 수 있어요!",
             Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.5),
             font_size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def add_slide_04(prs):
    """슬라이드 4: 헬멧 착용이 왜 중요한가"""
    sl = blank_slide(prs)
    header_bar(sl, "헬멧 착용이 왜 중요한가")
    # 왼쪽: 착용
    add_rect(sl, Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5), RGBColor(0xE3, 0xF2, 0xFD))
    add_text(sl, "✅ 헬멧 착용 시",
             Inches(0.7), Inches(1.4), Inches(5.4), Inches(0.6),
             font_size=20, bold=True, color=GREEN)
    items_on = ["• 충격 에너지 85% 흡수", "• 두개골 골절 위험 ↓↓", "• 뇌진탕 위험 크게 감소", "• 경미한 찰과상으로 끝남"]
    for j, item in enumerate(items_on):
        add_text(sl, item,
                 Inches(0.7), Inches(2.1 + j * 0.75), Inches(5.4), Inches(0.65),
                 font_size=18, color=DARK)
    # 오른쪽: 미착용
    add_rect(sl, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.5), RGBColor(0xFF, 0xEB, 0xEE))
    add_text(sl, "❌ 헬멧 미착용 시",
             Inches(7.2), Inches(1.4), Inches(5.4), Inches(0.6),
             font_size=20, bold=True, color=ORANGE)
    items_off = ["• 머리 직접 충격 흡수", "• 두개골 골절 위험 높음", "• 뇌손상·사망 가능성", "• 장기 후유증 위험"]
    for j, item in enumerate(items_off):
        add_text(sl, item,
                 Inches(7.2), Inches(2.1 + j * 0.75), Inches(5.4), Inches(0.65),
                 font_size=18, color=DARK)
    # 가운데 VS
    add_text(sl, "VS",
             Inches(6.0), Inches(2.8), Inches(1.33), Inches(1.0),
             font_size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def add_slide_05(prs):
    """슬라이드 5: 올바른 헬멧 착용법"""
    sl = blank_slide(prs)
    header_bar(sl, "올바른 헬멧 착용법")
    steps = [
        ("1단계", "수평으로 쓰기", "이마가 보이도록 앞이마 2~3cm 위에 수평으로 착용"),
        ("2단계", "귀 앞 V자 조절", "귀 아래를 기준으로 양쪽 끈이 V자가 되도록 조절"),
        ("3단계", "턱 끈 고정", "손가락 1~2개가 들어갈 정도로 턱 끈을 조인다"),
    ]
    colors = [BLUE, GREEN, ORANGE]
    for i, (step, title, desc) in enumerate(steps):
        y = Inches(1.4 + i * 1.8)
        add_rect(sl, Inches(0.5), y, Inches(1.5), Inches(1.4), colors[i])
        add_text(sl, step,
                 Inches(0.5), y + Inches(0.35),
                 Inches(1.5), Inches(0.65),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(2.2), y, Inches(10.5), Inches(1.4), LGRAY)
        add_text(sl, title,
                 Inches(2.4), y + Inches(0.05),
                 Inches(10.1), Inches(0.55),
                 font_size=20, bold=True, color=colors[i])
        add_text(sl, desc,
                 Inches(2.4), y + Inches(0.6),
                 Inches(10.1), Inches(0.65),
                 font_size=16, color=DARK)
    add_text(sl, "⚠️ 자전거를 탈 때는 항상 헬멧을 먼저 착용하세요!",
             Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
             font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def add_slide_06(prs):
    """슬라이드 6: 출발 전 점검 ABC"""
    sl = blank_slide(prs)
    header_bar(sl, "출발 전 점검 ABC", bg_color=GREEN)
    abc = [
        ("A", "Air\n(타이어 공기압)", "타이어를 손으로 눌러 단단한지 확인\n손가락으로 누를 때 거의 안 들어가면 OK", BLUE),
        ("B", "Brake\n(브레이크)", "앞·뒤 브레이크 레버를 각각 잡아보기\n확실히 멈추는지 테스트 후 출발", GREEN),
        ("C", "Chain\n(체인)", "체인이 빠지거나 끊어지지 않았는지 확인\n녹이나 오염이 심하면 교체 필요", ORANGE),
    ]
    for i, (letter, title, desc, color) in enumerate(abc):
        x = Inches(0.5 + i * 4.3)
        add_rect(sl, x, Inches(1.4), Inches(3.8), Inches(4.8), color)
        add_text(sl, letter,
                 x, Inches(1.5), Inches(3.8), Inches(1.4),
                 font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title,
                 x + Inches(0.1), Inches(2.9), Inches(3.6), Inches(0.9),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc,
                 x + Inches(0.1), Inches(3.8), Inches(3.6), Inches(1.2),
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


def add_slide_07(prs):
    """슬라이드 7: 도로 주행 규칙"""
    sl = blank_slide(prs)
    header_bar(sl, "도로 주행 규칙")
    rules = [
        ("🚲 자전거도로 우선", "자전거도로가 있으면 반드시 이용\n없으면 차도 우측 가장자리로 주행"),
        ("↕ 우측통행", "항상 도로 오른쪽으로 주행\n역주행은 매우 위험!"),
        ("🛑 신호 준수", "보행자 신호·차량 신호 모두 지키기\n빨간불에 절대 통과 금지"),
        ("👁 주변 확인", "교차로에서는 반드시 서서히 접근\n좌우 차량·보행자 확인 후 통과"),
    ]
    for i, (title, desc) in enumerate(rules):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.4 + row * 2.8)
        add_rect(sl, x, y, Inches(6.0), Inches(2.4), LGRAY)
        add_text(sl, title,
                 x + Inches(0.15), y + Inches(0.15),
                 Inches(5.7), Inches(0.7),
                 font_size=20, bold=True, color=BLUE)
        add_text(sl, desc,
                 x + Inches(0.15), y + Inches(0.85),
                 Inches(5.7), Inches(1.3),
                 font_size=16, color=DARK)


def add_slide_08(prs):
    """슬라이드 8: 야간·우천 시 안전 수칙"""
    sl = blank_slide(prs)
    header_bar(sl, "야간·우천 시 안전 수칙", bg_color=RGBColor(0x37, 0x47, 0x4F))
    tips = [
        ("💡 전조등·후미등 켜기", "야간에는 반드시 전조등(앞)과 후미등(뒤) 점등\n배터리 미리 충전 확인"),
        ("🔆 반사판·반사 조끼 착용", "야광 조끼나 반사 스티커로 내 위치를 알려요\n운전자 눈에 잘 띄어야 안전"),
        ("🌧 우천 시 속도 줄이기", "빗길은 제동 거리가 2배 이상 길어짐\n평소보다 훨씬 느리게 달리기"),
        ("🚫 우천 시 우산 들고 타기 금지", "한 손 운전은 균형을 잃기 쉬움\n비 올 때는 우비+헬멧 세트로 착용"),
    ]
    for i, (title, desc) in enumerate(tips):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.4 + row * 2.8)
        add_rect(sl, x, y, Inches(6.0), Inches(2.4), RGBColor(0x37, 0x47, 0x4F))
        add_text(sl, title,
                 x + Inches(0.15), y + Inches(0.15),
                 Inches(5.7), Inches(0.7),
                 font_size=18, bold=True, color=WHITE)
        add_text(sl, desc,
                 x + Inches(0.15), y + Inches(0.85),
                 Inches(5.7), Inches(1.3),
                 font_size=15, color=RGBColor(0xE0, 0xE0, 0xE0))


def add_slide_09(prs):
    """슬라이드 9: 핸드폰·이어폰 사용 금지"""
    sl = blank_slide(prs)
    header_bar(sl, "핸드폰·이어폰 사용 금지", bg_color=ORANGE)
    # 큰 금지 아이콘 영역
    add_rect(sl, Inches(0.5), Inches(1.4), Inches(4.0), Inches(4.5), RGBColor(0xFF, 0xEB, 0xEE))
    add_text(sl, "📵",
             Inches(0.5), Inches(1.8), Inches(4.0), Inches(2.0),
             font_size=72, align=PP_ALIGN.CENTER)
    add_text(sl, "자전거 탑승 중\n핸드폰·이어폰\n절대 금지!",
             Inches(0.5), Inches(3.8), Inches(4.0), Inches(1.8),
             font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    # 오른쪽 이유 설명
    reasons = [
        ("⚠️ 시야 차단", "핸드폰 화면을 보는 순간 전방 주시 불가"),
        ("⚠️ 반응 속도 저하", "위험 상황에 1~2초 늦게 반응 → 사고"),
        ("⚠️ 소리 차단", "이어폰 착용 시 경적·경고음을 못 들음"),
        ("📋 법적 처벌", "도로교통법 위반 시 범칙금 부과 가능"),
    ]
    for i, (title, desc) in enumerate(reasons):
        y = Inches(1.4 + i * 1.1)
        add_rect(sl, Inches(5.0), y, Inches(7.8), Inches(0.95), LGRAY)
        add_text(sl, title,
                 Inches(5.15), y + Inches(0.05),
                 Inches(3.5), Inches(0.4),
                 font_size=16, bold=True, color=ORANGE)
        add_text(sl, desc,
                 Inches(5.15), y + Inches(0.45),
                 Inches(7.5), Inches(0.4),
                 font_size=14, color=DARK)


def build_ppt():
    """Generate the bicycle safety education presentation and save it."""
    prs = new_prs()
    add_slide_01(prs)
    add_slide_02(prs)
    add_slide_03(prs)
    add_slide_04(prs)
    add_slide_05(prs)
    add_slide_06(prs)
    add_slide_07(prs)
    add_slide_08(prs)
    add_slide_09(prs)
    prs.save("자전거_안전_교육.pptx")
    print("저장 완료: 자전거_안전_교육.pptx")


if __name__ == "__main__":
    build_ppt()
